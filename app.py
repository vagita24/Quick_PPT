from flask import Flask,render_template,url_for,request,send_file
import os
from pptx import Presentation
import subprocess

app = Flask(__name__)

@app.route('/')             #mapping of url (in '') to return statement ,'@' known as decorator
def home():
    return render_template('home.html')

@app.route('/return-file/')             #mapping of url (in '') to return statement ,'@' known as decorator
def return_file():
    return send_file('/home/sahil/Desktop/text_summarize/static/test.pptx')

@app.route('/file-downloads/')             #mapping of url (in '') to return statement ,'@' known as decorator
def file_downloads():
    return render_template('download.html')

@app.route('/predict',methods=['POST'])
def predict():
    if request.method == 'POST':
        comment = request.form['comment']    
        import bs4 as bs  
        import urllib.request  
        import re
        import nltk
        nltk.download('punkt')
        nltk.download('stopwords')
        import heapq



        scraped_data = urllib.request.urlopen(comment)  
        article = scraped_data.read()

        parsed_article = bs.BeautifulSoup(article,'lxml')

        title = parsed_article.find_all('title')
        new = str(title[0])
        new = re.sub('[^a-zA-Z]',' ',new)
        new = re.sub(r'\s+',' ',new)
        new = new.replace('title','')

        paragraphs = parsed_article.find_all('p')

        article_text = ""
        current_dir='/home/sahil/Desktop/text_summarize'
        prs = Presentation(current_dir+'/template_sample.pptx')
        subprocess.run(['rm','-rf',f'{current_dir}/test.pptx'])
        for p in paragraphs:
            article_text = p.text
            article_text = re.sub(r'\s+', ' ', article_text)
            formatted_article_text = re.sub('[^a-zA-Z]', ' ', article_text )
            formatted_article_text = re.sub(r'\s+', ' ', formatted_article_text)
            #converting text to sentence list, thus we use original article(article_text) coz it contains '.'' 
            sentence_list = nltk.sent_tokenize(article_text)  

            if (sentence_list):
            #Find weighted frequency of occurances
                stopwords = nltk.corpus.stopwords.words('english')

                word_frequencies = {}  
                for word in nltk.word_tokenize(formatted_article_text):  
                    if word not in stopwords:
                        if word not in word_frequencies.keys():
                            word_frequencies[word] = 1
                        else:
                            word_frequencies[word] += 1


                # print(word_frequencies)

                #find max accuracy word and normalize others

                maximum_frequncy = max(word_frequencies.values())

                for word in word_frequencies.keys():  
                    word_frequencies[word] = (word_frequencies[word]/maximum_frequncy)

                #print(word_frequencies)

                #calculate sentence score

                sentence_scores = {}  
                for sent in sentence_list:  
                    for word in nltk.word_tokenize(sent.lower()):
                        if word in word_frequencies.keys():
                            if len(sent.split(' ')) < 30:
                                if sent not in sentence_scores.keys():
                                    sentence_scores[sent] = word_frequencies[word]
                                else:
                                    sentence_scores[sent] += word_frequencies[word]

                #print(sentence_scores)

                #Predicting top 7 sentences

                summary_sentences = heapq.nlargest(4, sentence_scores, key=sentence_scores.get)
                
            #summary = ' '.join(summary_sentences)  

    

                if(len(summary_sentences)>=3):
                    title_slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(title_slide_layout)
                    #title = slide.shapes.title
                    #subtitle = slide.placeholders[1]

                    title.text = new
                    shapes=slide.shapes
                    body_shape=shapes.placeholders[1]
                    tf=body_shape.text_frame
                    #SLD_LAYOUT_TITLE_AND_CONTENT = 1
                    #slide_layout= prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
                    #slide = prs.slides.add_slide(slide_layout)
                    for p in summary_sentences:
                        para=tf.add_paragraph()
                        para.text=p
                        para.level=1
                


        prs.save('/home/sahil/Desktop/text_summarize/test.pptx')
    return render_template('download.html')

if __name__ == '__main__':      #this is our main file and we have to run this directly
    app.run(debug=True)
